from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_SLIDE = RGBColor(0x1e, 0x1e, 0x36)
ACCENT = RGBColor(0x00, 0xd2, 0xff)
ACCENT2 = RGBColor(0x7c, 0x3a, 0xed)
GREEN = RGBColor(0x50, 0xfa, 0x7b)
ORANGE = RGBColor(0xff, 0xb7, 0x00)
RED = RGBColor(0xff, 0x55, 0x55)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0xaa, 0xaa, 0xcc)
DARK_CARD = RGBColor(0x28, 0x28, 0x44)

def set_slide_bg(slide, color=BG_SLIDE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_card(slide, left, top, width, height, title, content, title_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(4)

# === SLIDE 1: Title ===
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 1.5, 11.7, 1.2, "RFID Manager", font_size=52, color=WHITE, bold=True)
add_text_box(slide, 0.8, 2.8, 11.7, 0.8, "Test Setup — Overview", font_size=32, color=ACCENT, bold=False)

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.7), Inches(3), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

add_text_box(slide, 0.8, 4.0, 11.7, 0.5, "MQTT Infrastructure • NFC Scanning • End-to-End Testing", font_size=18, color=GRAY)
add_text_box(slide, 0.8, 5.5, 11.7, 0.5, "2026-06-14", font_size=16, color=GRAY)

# === SLIDE 2: System Overview ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "System Overview", font_size=36, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.0, 12, 0.4, "Full architecture from NFC scan to data persistence", font_size=16, color=GRAY)

# Architecture diagram - top row: 3 blocks
add_card(slide, 0.5, 1.6, 3.8, 2.5,
    "📱 Android App (RFID Manager)",
    "• NFC Reader/Writer\n• MQTT Client (Paho)\n• Room DB (local cache)\n• Scan → Publish flow\n• IP: 192.168.50.112\n• WiFi-connected")

add_card(slide, 4.7, 1.6, 3.8, 2.5,
    "🌐 MQTT Broker (Mosquitto)",
    "• Docker container\n• eclipse-mosquitto:latest\n• Port 1883 (TCP)\n• allow_anonymous: true\n• IP: 192.168.50.128\n• Mosquitto 2.1.2")

add_card(slide, 8.9, 1.6, 3.8, 2.5,
    "💾 Subscriber / Backend",
    "• Python subscriber script\n• SQLite persistence\n• Filters: rfidmanager/+/telemetry\n• Color-coded terminal output\n• Statistics on exit\n• Localhost connection")

# Bottom row: data flow arrow placeholders
add_card(slide, 0.5, 4.5, 5.8, 2.5,
    "🔁 Data Flow",
    "1. NFC-tag scanned → UID extracted\n2. JSON payload constructed\n3. Published to rfidmanager/<uid>/telemetry\n4. Broker routes to all subscribers\n5. Python subscriber logs to SQLite\n6. App stores locally in Room DB")

add_card(slide, 6.7, 4.5, 5.8, 2.5,
    "🧪 Test Configuration",
    "• Broker: 192.168.50.128:1883\n• Topic: rfidmanager/<uid>/telemetry\n• QoS: 1 (at least once)\n• Keep-alive: 30s\n• cleanSession: true\n• JSON format payload")

# === SLIDE 3: MQTT Infrastructure ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "MQTT Infrastructure", font_size=36, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.0, 12, 0.4, "Publish/Subscribe pattern — decoupled architecture", font_size=16, color=GRAY)

# Left: Publish/subscribe explanation
add_card(slide, 0.5, 1.6, 6.0, 5.0,
    "Publish / Subscribe Model",
    "┌─────────────┐     ┌─────────────┐\n"
    "│  Publisher   │     │  Publisher  │\n"
    "│  (App)       │     │  (Simulator)│\n"
    "└──────┬──────┘     └──────┬──────┘\n"
    "       │                   │\n"
    "       ▼                   ▼\n"
    "┌─────────────────────────────────┐\n"
    "│      MQTT BROKER (Mosquitto)     │\n"
    "│   192.168.50.128:1883            │\n"
    "└────────┬────────────────┬────────┘\n"
    "         │                │\n"
    "         ▼                ▼\n"
    "┌─────────────┐     ┌─────────────┐\n"
    "│  Subscriber  │     │  Subscriber │\n"
    "│  (Python)    │     │  (MQTT Exp) │\n"
    "│  → SQLite    │     │  → GUI view │\n"
    "└─────────────┘     └─────────────┘\n\n"
    "Key points:\n"
    "• Publishers don't know subscribers\n"
    "• Broker handles all routing\n"
    "• Any number of pub/sub on each topic\n"
    "• Wildcards enable flexible subscriptions")

# Right: Topic structure
add_card(slide, 6.8, 1.6, 5.8, 5.0,
    "Topic Structure",
    "rfidmanager / <uid> / telemetry\n"
    "    ↑          ↑          ↑\n"
    "  our app    tag ID    data type\n\n"
    "Examples:\n"
    "• rfidmanager/047B1234/telemetry\n"
    "• rfidmanager/0A1B2C3D/telemetry\n\n"
    "Wildcards:\n"
    "• rfidmanager/#  — all messages\n"
    "• rfidmanager/+/telemetry  — all UIDs\n"
    "• rfidmanager/047B1234/+  — all types\n\n"
    "JSON Payload:\n"
    '{"uid":"047B1234","type":"rfid",\n'
    ' "timestamp":1718300000}\n\n'
    "QoS Levels:\n"
    "• 0: fire-and-forget\n"
    "• 1: at least once (our choice)\n"
    "• 2: exactly once (overkill)")

# === SLIDE 4: Network Topology ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "Network Topology", font_size=36, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.0, 12, 0.4, "Home network: 192.168.50.0/24", font_size=16, color=GRAY)

# Network diagram as cards
add_card(slide, 0.5, 1.6, 4.0, 2.2,
    "📱 Phone (WiFi)",
    "• Samsung Galaxy Note 10\n• IP: 192.168.50.112\n• Paho MQTT Client\n• clientId: rfid-android-client\n• Publishes on scan")

add_card(slide, 4.8, 1.6, 4.0, 2.2,
    "💻 PC — Docker Host",
    "• IP: 192.168.50.128\n• Mosquitto container\n  (port 1883)\n• Python subscriber\n• MQTT Explorer\n• All on same machine")

add_card(slide, 9.1, 1.6, 3.6, 2.2,
    "🌍 External Tools",
    "• mosquitto_pub/sub\n  (standalone or via docker)\n• MQTT Explorer (GUI)\n• Wireshark (packet analysis)\n• netcat (raw TCP test)")

# Bottom: connection details
add_card(slide, 0.5, 4.2, 12.0, 3.0,
    "Connection Matrix",
    "┌──────────────────────────────────────────────────────────────────────────────┐\n"
    "│ Component         │ IP/Host          │ Connects To       │ Port │ Protocol │\n"
    "├──────────────────────────────────────────────────────────────────────────────┤\n"
    "│ Phone (App)       │ 192.168.50.112   │ Broker (WiFi)     │ 1883 │ MQTT     │\n"
    "│ Broker (Container) │ 192.168.50.128  │ — (accepts all)   │ 1883 │ TCP      │\n"
    "│ Python Subscriber │ localhost        │ Broker (local)    │ 1883 │ MQTT     │\n"
    "│ MQTT Explorer     │ 192.168.50.128   │ Broker            │ 1883 │ MQTT     │\n"
    "│ mosquitto CLI     │ localhost/192.168.50.128 │ Broker    │ 1883 │ MQTT     │\n"
    "│ Wireshark         │ any              │ Network interface │ —    │ MQTT ftr │\n"
    "└──────────────────────────────────────────────────────────────────────────────┘")

# === SLIDE 5: Test Tools ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "Test Tools & Commands", font_size=36, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.0, 12, 0.4, "Three ways to test — CLI, Docker, or GUI", font_size=16, color=GRAY)

# Tool cards
add_card(slide, 0.5, 1.6, 6.0, 2.5,
    "🔧 Method 1: Standalone mosquitto CLI",
    "Install:  sudo pacman -S mosquitto\n\n"
    "Listen on all topics:\n"
    "  mosquitto_sub -h 192.168.50.128 -p 1883 \\\n"
    "    -t \"rfidmanager/#\" -v\n\n"
    "Publish test message:\n"
    "  mosquitto_pub -h 192.168.50.128 -p 1883 \\\n"
    "    -t \"rfidmanager/test/telemetry\" \\\n"
    "    -m '{\"uid\":\"test\"}'")

add_card(slide, 6.8, 1.6, 6.0, 2.5,
    "🐳 Method 2: Via Docker container",
    "No installation needed — broker already has mosquitto tools.\n\n"
    "Listen:\n"
    "  docker exec rfid-mqtt-test mosquitto_sub \\\n"
    "    -h localhost -p 1883 -t \"rfidmanager/#\" -v\n\n"
    "Publish:\n"
    "  docker exec rfid-mqtt-test mosquitto_pub \\\n"
    "    -h localhost -p 1883 \\\n"
    "    -t \"rfidmanager/test/telemetry\" \\\n"
    "    -m '{\"type\":\"test\"}'")

add_card(slide, 0.5, 4.4, 6.0, 2.5,
    "🖥️ Method 3: MQTT Explorer (GUI)",
    "Download: github.com/thomasnordquist/MQTT-Explorer\n\n"
    "Connection settings:\n"
    "  Host: 192.168.50.128\n"
    "  Port: 1883\n"
    "  SSL/TLS: Off\n"
    "  Auth: None\n\n"
    "Subscribe to rfidmanager/# for real-time view")

add_card(slide, 6.8, 4.4, 6.0, 2.5,
    "🐍 Python Test Scripts",
    "Location: ~/rfid-manager/test/fas2-mqtt/mqtt/\n\n"
    "Subscriber (SQLite logging):\n"
    "  ./test_subscriber_persist.py\n"
    "  ./test_subscriber_persist.py --uid 047B\n\n"
    "Simulator (simulate phone):\n"
    "  ./simulate_mobile_publish.py\n\n"
    "Output: ~/rfid-manager/data/rfid_readings.db")

# === SLIDE 6: Broker Administration ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "Broker Administration", font_size=36, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.0, 12, 0.4, "Docker-based Mosquitto management", font_size=16, color=GRAY)

add_card(slide, 0.5, 1.6, 6.0, 5.0,
    "Docker Commands",
    "┌─────────────────────────────────────────────────────────────┐\n"
    "│ Command                           │ Description             │\n"
    "├─────────────────────────────────────────────────────────────┤\n"
    "│ docker run ... eclipse-mosquitto  │ Start broker (--rm!)   │\n"
    "│ docker start rfid-mqtt-test       │ Start existing cont.   │\n"
    "│ docker stop rfid-mqtt-test        │ Stop broker           │\n"
    "│ docker restart rfid-mqtt-test     │ Restart broker        │\n"
    "│ docker ps --filter name=rfid-...  │ Check if running      │\n"
    "│ docker logs rfid-mqtt-test        │ View logs             │\n"
    "│ docker logs -f rfid-mqtt-test     │ Follow logs live      │\n"
    "│ docker exec -it rfid-mqtt-test sh  │ Shell inside container │\n"
    "│ docker inspect rfid-mqtt-test     │ Detailed info         │\n"
    "└─────────────────────────────────────────────────────────────┘\n\n"
    "⚠️ The container uses --rm: it is deleted on stop.\n"
    "   You must run 'docker run' again, not 'docker start'.\n"
    "   Omit --rm if you want persistent container lifecycle.")

add_card(slide, 6.8, 1.6, 5.8, 5.0,
    "mosquitto.conf",
    "File: ~/rfid-manager/test/fas2-mqtt/mqtt/mosquitto.conf\n\n"
    "Config:\n"
    "  listener 1883              # MQTT port\n"
    "  allow_anonymous true       # No auth required\n"
    "  persistence true           # Data saved to disk\n"
    "  persistence_location       # /mosquitto/data/\n"
    "  log_dest stdout           # Logs to console\n\n"
    "Not configured (defaults):\n"
    "  • max_connections: unlimited\n"
    "  • keepalive_interval: 60s\n"
    "  • allow_zero_length_clientid: true\n\n"
    "Logs show:\n"
    "  • New connections (IP + port)\n"
    "  • Client IDs\n"
    "  • Disconnections")

# === SLIDE 7: Test Flow ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "End-to-End Test Flow", font_size=36, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.0, 12, 0.4, "From NFC scan to data in SQLite — complete chain", font_size=16, color=GRAY)

# Flow steps as horizontal cards
flow_steps = [
    ("1️⃣", "NFC Scan", "Phone reads\nRFID tag\n→ extracts UID"),
    ("2️⃣", "Build Payload", "JSON constructed:\nuid, type,\ntimestamp"),
    ("3️⃣", "Publish MQTT", "→ rfidmanager/\n<uid>/telemetry\nQoS 1"),
    ("4️⃣", "Broker Routes", "Mosquitto receives\nand forwards\nto subscribers"),
    ("5️⃣", "Subscriber", "Python script\nreceives →\nparses JSON"),
    ("6️⃣", "SQLite Store", "Data persisted\nin rfid_readings.db\nfor analysis"),
]

for i, (emoji, title, desc) in enumerate(flow_steps):
    left = 0.5 + i * 2.1
    add_card(slide, left, 1.6, 1.9, 2.8, f"{emoji} {title}", desc, ACCENT)

# Verification methods
add_card(slide, 0.5, 4.8, 12.0, 2.3,
    "✅ How to Verify",
    "┌─────────────────────────────────────────────────────────────────────────────────┐\n"
    "│ Method                    │ What to check                              │ Pass    │\n"
    "├─────────────────────────────────────────────────────────────────────────────────┤\n"
    "│ Terminal: python subscr.  │ See incoming messages with UID and timestamp │ ✅      │\n"
    "│ Terminal: mosquitto_sub   │ -v flag shows exact topic + payload          │ ✅      │\n"
    "│ GUI: MQTT Explorer        │ Real-time tree of all topics                │ ✅      │\n"
    "│ App: Settings             │ Broker status indicator                      │ ✅      │\n"
    "│ SQLite: rfid_readings.db  │ Query: SELECT * FROM readings;               │ ✅      │\n"
    "│ App: Readings screen      │ Scan appears in list                        │ ✅      │\n"
    "└─────────────────────────────────────────────────────────────────────────────────┘")

# === SLIDE 8: Quick Reference ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "Quick Reference", font_size=36, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.0, 12, 0.4, "Essential commands and configurations", font_size=16, color=GRAY)

add_card(slide, 0.5, 1.6, 4.0, 5.0,
    "📡 Connection Details",
    "┌──────────────────────────────┐\n"
    "│ Parameter    │ Value         │\n"
    "├──────────────────────────────┤\n"
    "│ Broker IP    │ 192.168.50.128│\n"
    "│ Port         │ 1883          │\n"
    "│ Protocol     │ MQTT 3.1.1   │\n"
    "│ Topic root   │ rfidmanager/  │\n"
    "│ QoS          │ 1             │\n"
    "│ Auth         │ None (anon)   │\n"
    "│ TLS          │ No            │\n"
    "│ Keep-alive   │ 30s           │\n"
    "│ Clean sess   │ true          │\n"
    "└──────────────────────────────┘")

add_card(slide, 4.8, 1.6, 4.0, 5.0,
    "🗂️ Important Files",
    "┌──────────────────────────────────────┐\n"
    "│ Path │ Purpose                      │\n"
    "├──────────────────────────────────────┤\n"
    "│ mosquitto.conf                      │\n"
    "│ (test/fas2-mqtt/mqtt/)               │\n"
    "│                                     │\n"
    "│ test_subscriber_persist.py          │\n"
    "│ (test/fas2-mqtt/mqtt/)               │\n"
    "│                                     │\n"
    "│ simulate_mobile_publish.py          │\n"
    "│ (test/fas2-mqtt/mqtt/)               │\n"
    "│                                     │\n"
    "│ rfid_readings.db                    │\n"
    "│ (data/)                              │\n"
    "│                                     │\n"
    "│ MQTT-Infrastruktur.md               │\n"
    "│ (wiki/)                              │\n"
    "└──────────────────────────────────────┘")

add_card(slide, 9.1, 1.6, 3.6, 5.0,
    "🧪 One-Liners",
    "Verify broker is up:\n"
    "  docker ps | grep mqtt\n\n"
    "Listen to everything:\n"
    "  docker exec rfid-mqtt-test \\\n"
    "    mosquitto_sub \\\n"
    "    -h localhost -p 1883 \\\n"
    "    -t \"rfidmanager/#\" -v\n\n"
    "Publish test:\n"
    "  docker exec rfid-mqtt-test \\\n"
    "    mosquitto_pub -h localhost \\\n"
    "    -p 1883 -t \"test\" -m \"ping\"\n\n"
    "View logs:\n"
    "  docker logs rfid-mqtt-test")

# === SLIDE 9: End ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 2.0, 11.7, 1.0, "Questions?", font_size=48, color=WHITE, bold=True)
add_text_box(slide, 0.8, 3.2, 11.7, 0.6, "RFID Manager — MQTT Test Infrastructure", font_size=24, color=ACCENT)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.0), Inches(3), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

add_text_box(slide, 0.8, 4.3, 11.7, 0.5, "Joakim • 2026-06-14", font_size=18, color=GRAY)

# Save
output_path = "/home/joakim/llm-wiki/RFID-Manager-Test-Setup.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
