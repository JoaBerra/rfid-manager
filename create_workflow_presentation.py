from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_SLIDE = RGBColor(0x1e, 0x1e, 0x36)
ACCENT = RGBColor(0x00, 0xd2, 0xff)
ACCENT2 = RGBColor(0x7c, 0x3a, 0xed)
GREEN = RGBColor(0x50, 0xfa, 0x7b)
ORANGE = RGBColor(0xff, 0xb7, 0x00)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0xaa, 0xaa, 0xcc)
DARK_CARD = RGBColor(0x28, 0x28, 0x44)

def set_bg(slide, color=BG_SLIDE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tbox(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb

def card(slide, l, t, w, h, title, content, title_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.space_after = Pt(8)
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(4)

# === SLIDE 1: Title ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
tbox(s, 0.8, 1.5, 11.7, 1.2, "RFID Manager", size=52, color=WHITE, bold=True)
tbox(s, 0.8, 2.8, 11.7, 0.8, "Arbetssätt", size=32, color=ACCENT)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.7), Inches(3), Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
tbox(s, 0.8, 4.0, 11.7, 0.5, "Utvecklingsflöde • Verktyg • Roller • AI-assistans", size=18, color=GRAY)
tbox(s, 0.8, 5.5, 11.7, 0.5, "2026-06-14", size=16, color=GRAY)

# === SLIDE 2: Overview ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Vårt arbetssätt — översikt", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "AI-assisterad utveckling med tydlig struktur och iteration", size=16, color=GRAY)

card(s, 0.5, 1.6, 3.8, 2.5,
    "🎯 Målstyrt",
    "• Varje fas har tydliga\n  acceptanskriterier\n• Kund (projektledare)\n  godkänner innan\n  nästa steg\n• Prioritering via\n  Produkt-Roadmap")

card(s, 4.7, 1.6, 3.8, 2.5,
    "🔄 Iterativt",
    "• Korta cykler:\n  bygg → testa →\n  godkänn\n• AI-assistenten\n  utför, kunden\n  granskar\n• Kontinuerlig\n  förbättring")

card(s, 8.9, 1.6, 3.8, 2.5,
    "📝 Dokumenterat",
    "• Allt finns i wikin\n  (Karpathy-pattern)\n• Beslut loggas\n• Återanvändbara\n  artefakter\n• Snabb onboarding")

card(s, 0.5, 4.5, 5.8, 2.5,
    "🤖 AI-assistans (OpenCode)",
    "• Chat-interface: instruktioner → kod\n• AI utför implementation, tester, dokumentation\n• Kunden granskar och godkänner\n• Transparent — allt visas i terminalen\n• AI-agenter för parallella uppgifter")

card(s, 6.7, 4.5, 5.8, 2.5,
    "📊 Tre nivåer av planering",
    "• Roadmap (Fil: Produkt-Roadmap.md)\n  — Långsiktig vision, faser 1-6+\n• Kanban (Fil: Kanban.md)\n  — Veckovis taktik, att göra/pågår/klart\n• Session log (dokumenteras löpande)\n  — Vad gjordes, beslut, nästa steg")

# === SLIDE 3: Verktyg ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Verktyg", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Helst open source, minimal overhead", size=16, color=GRAY)

tools_info = [
    ("🧠 OpenCode", "AI-assistent", "CLI-baserad AI som utför\nkod, tester, dokumentation.\nIntegrerad med terminalen."),
    ("📓 Obsidian", "Wiki/Kunskapsbas", "All dokumentation i markdown\nmed Karpathy-pattern:\nstort index → djupdykning i delområden."),
    ("📋 Kanban (Obsidian)", "Task management", "Inbyggd i wikin via plugin.\nAtt göra → Pågår → Klart.\nLänkar till wiki-sidor."),
    ("💻 Neovim", "Editor", "Primär kod-editor.\nAI-terminal via OpenCode."),
    ("🐙 GitHub", "Versionshantering", "Två repos:\nrfid-manager (wiki)\nrfid-manager-android (kod)."),
    ("🐳 Docker", "Miljö", "Mosquitto broker,\ntestverktyg, isolering."),
]

for i, (emoji, title, desc) in enumerate(tools_info):
    col = i % 3
    row = i // 3
    l = 0.5 + col * 4.2
    t = 1.6 + row * 2.8
    card(s, l, t, 3.8, 2.4, f"{emoji} {title}", desc)

# === SLIDE 4: Utvecklingsflöde ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Utvecklingsflöde", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Från behov till godkänd funktion — steg för steg", size=16, color=GRAY)

flow_steps = [
    ("1", "Behov", "Kund/PM identifierar\nen ny funktion eller\nförbättring"),
    ("2", "Planering", "Läggs till i Kanban\nAtt göra.\nPrioriteras."),
    ("3", "Implementation", "AI implementerar\nenligt krav och\naccepanskriterier"),
    ("4", "Demo", "Kunden granskar\nlive: testkör,\ninspekterar"),
    ("5", "Godkännande", "Kund godkänner →\nflyttas till Klart\ni Kanban"),
    ("6", "Dokumentation", "Wiki uppdateras:\nbeslut, konfig,\nkommandon"),
]

for i, (num, title, desc) in enumerate(flow_steps):
    l = 0.5 + i * 2.1
    card(s, l, 1.6, 1.9, 2.5, f"{num}. {title}", desc, ACCENT)

# Arrows between steps (text-based in card below)
card(s, 0.5, 4.5, 12.0, 2.7,
    "🔁 Iterationsloop",
    "┌─────────────────────────────────────────────────────────────────────────────────────┐\n"
    "│                                                                                     │\n"
    "│   Behov → Kanban (Att göra) → AI implementerar → Kund testar → Godkänt?             │\n"
    "│                                                      │                               │\n"
    "│                                              ┌───────┴────────┐                      │\n"
    "│                                              │  Ja → Klart +  │                      │\n"
    "│                                              │  Dokumentation  │                      │\n"
    "│                                              │  Nej → Justera  │                      │\n"
    "│                                              └────────────────┘                      │\n"
    "│                                                                                     │\n"
    "│  Aldrig klarmarkerad före kund-godkännande. AI-assistenten håller boarden aktuell.  │\n"
    "└─────────────────────────────────────────────────────────────────────────────────────┘")

# === SLIDE 5: Roller ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Roller & Ansvar", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Tydlig rollfördelning i AI-assisterad utveckling", size=16, color=GRAY)

card(s, 0.5, 1.6, 5.8, 2.5,
    "👤 Kund / Projektledare (Du)",
    "• Säger VAD — definierar behov och acceptanskriterier\n• Prioriterar — bestämmer vad som görs härnäst\n• Granskar — testar och godkänner varje funktion\n• Äger produkten — slutgiltiga beslut\n• Signerar acceptans i Kundrelationer-och-Acceptans.md")

card(s, 6.8, 1.6, 5.8, 2.5,
    "🤖 AI-assistent (OpenCode)",
    "• Säger HUR — föreslår implementation, verktyg, struktur\n• Utför — skriver kod, tester, dokumentation\n• Förklarar — pedagogiskt, med tabeller och diagram\n• Kommer ihåg — session context, beslut, preferenser\n• Håller Kanban uppdaterad")

card(s, 0.5, 4.5, 5.8, 2.5,
    "📋 Kommunikationsmönster",
    "• Kund: \"Gör X\" → AI utför → \"Vad tycker du?\"\n• AI: \"Här är resultatet...\" + fråga om nästa steg\n• Uppföljningsfrågor uppmuntras\n• All dialog sparas i kontexten under sessionen\n• Nya sessioner börjar med: \"Vad gjorde vi sist?\"")

card(s, 6.8, 4.5, 5.8, 2.5,
    "✅ Godkännandeprocess",
    "• Efter implementation: AI visar resultat\n• Kund testar själv eller granskar\n• Om godkänt: markeras i Kanban + sign-off\n• Om inte: justering + ny granskning\n• Alla godkännanden dokumenteras:\n  \"✅ Godkänt 2026-06-14\" med testfallsreferens")

# === SLIDE 6: Wiki ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Wiki-struktur (Karpathy-pattern)", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Ett index → många djupdykningar — skalbar kunskapsbas", size=16, color=GRAY)

card(s, 0.5, 1.6, 5.8, 5.0,
    "📄 index.md (navet)",
    "• Huvudingång till all dokumentation\n• Kategoriserade sektioner\n• Länkar till specialiserade sidor\n• Underhålls: döda länkar tas bort\n• Alla sidor måste finnas\n\n"
    "Exempelstruktur:\n"
    "  ┌─────────────────────────────────┐\n"
    "  │ index.md                        │\n"
    "  ├─────────────────────────────────┤\n"
    "  │ • Om project                    │\n"
    "  │ • Installation                  │\n"
    "  │ • MQTT-Infrastruktur            │\n"
    "  │ • App-Architecture              │\n"
    "  │ • Kanban                        │\n"
    "  │ • Produkt-Roadmap               │\n"
    "  │ • Kundrelationer-och-Acceptans  │\n"
    "  │ • ...                           │\n"
    "  └─────────────────────────────────┘")

card(s, 6.8, 1.6, 5.8, 5.0,
    "📑 Specialiserade sidor",
    "• Varje ämne har sin egen sida\n• Innehåller: beskrivning, konfig,\n  kommandon, exempel\n• Länkar mellan sidor (wiki-länkar)\n• Kanban innehåller bara status,\n  detaljer i separata sidor\n\n"
    "Exempel (MQTT-Infrastruktur.md):\n"
    "  1. MQTT-protokollet\n"
    "  2. Vår broker (Mosquitto)\n"
    "  3. Topologi och nätverk\n"
    "  4. Topics\n"
    "  5. Meddelandeformat\n"
    "  6. QoS\n"
    "  7. Återanslutning\n"
    "  8. Säkerhet\n"
    "  9. Verktyg\n"
    " 10. Flöde\n"
    " 11. Begränsningar\n"
    " 12. Referenser\n\n"
    "Format: Markdown med rubriker,\n"
    "tabeller, kodblock, ASCII-diagram")

# === SLIDE 7: Kanban ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Kanban Board", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Inbyggd i Obsidian — enkel men kraftfull", size=16, color=GRAY)

card(s, 0.5, 1.6, 3.8, 5.0,
    "🔜 Att göra (To Do)",
    "• Planerade uppgifter\n• Prioriteras av kund\n• Innehåller hela faser\n  nedbrutna i tasks\n\n"
    "Exempel:\n"
    "  [ ] Fas-100 MQTT-\n      infrastruktur\n    [x] Wiki-sida skapad\n    [ ] 1. MQTT-grunder\n    [ ] 2. Vår broker\n    [ ] 3. Topologi\n    ...\n\n"
    "Checkboxar för spårbarhet")

card(s, 4.7, 1.6, 3.8, 5.0,
    "🔄 Pågår (In Progress)",
    "• Aktuellt arbetsområde\n• Max 1-2 saker i taget\n• AI-assistenten flyttar\n  hit när arbete påbörjas\n\n"
    "Flöde:\n"
    "  Att göra → Pågår → Klart\n"
    "       ↑                  ↓\n"
    "       └── Justering ────┘")

card(s, 8.9, 1.6, 3.8, 5.0,
    "✅ Klart (Done)",
    "• Godkända funktioner\n• Datum för godkännande\n• Testfallsreferens\n• Sign-off i Kundrelationer\n\n"
    "Exempel:\n"
    "  [x] 6.1 MQTT-konfig i UI\n  ✅ Godkänt 2026-06-13\n\n"
    "  [x] BUG-002 Write\n  TC-SCAN-006, TC-E2E-003\n  ✅ Godkänt 2026-06-13\n\n"
    "Regel: aldrig klarmarkerad\nföre kund-godkännande")

# === SLIDE 8: Artefakter ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Artefakter & Dokument", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Allt som produceras under utvecklingen", size=16, color=GRAY)

card(s, 0.5, 1.6, 3.8, 2.5,
    "📱 Android App",
    "• Kotlin + Jetpack Compose\n• Paho MQTT-klient\n• Room-databas (SQLite)\n• NFC Reader/Writer\n• Material 3 UI\n• GitHub: rfid-manager-\n  android")

card(s, 4.7, 1.6, 3.8, 2.5,
    "📚 Wiki",
    "• index.md — navigation\n• MQTT-Infrastruktur.md\n• App-Architecture.md\n• Produkt-Roadmap.md\n• Kanban.md\n• Kundrelationer.md\n• Testplan.md\n• GitHub: rfid-manager")

card(s, 8.9, 1.6, 3.8, 2.5,
    "🧪 Test & Verktyg",
    "• Python subscriber (+ SQLite)\n• Python simulator\n• Testplan med 22 testfall\n• Användarmanual (PDF)\n• Presentations-PDFer\n• Arkitekturdiagram")

card(s, 0.5, 4.5, 3.8, 2.5,
    "📦 Releases",
    "• v1.0 (GitHub Release)\n• APK: app-release.apk\n• Release notes per fas\n• Debug-signerad\n• R8-minifierad\n• 2.2 MB nedladdning")

card(s, 4.7, 4.5, 3.8, 2.5,
    "📊 Presentations-PDFer",
    "• RFID-Manager-Test-Setup.pdf\n  — teknisk översikt\n• RFID-Manager-Arbetssatt.pdf\n  — processbeskrivning\n• Skapas med python-pptx\n  + LibreOffice\n• 16:9 format, mörkt tema")

card(s, 8.9, 4.5, 3.8, 2.5,
    "🗺️ Framtida artefakter",
    "• Användarmanual (PDF, klart)\n• Nomenklaturdokument (topics)\n• Broker quick reference (PDF)\n• Docker-compose för CI\n• Play Store-lansering\n• TLS-certifikatsguide")

# === SLIDE 9: Processregler ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Processregler", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Spelregler som håller arbetet fokuserat", size=16, color=GRAY)

rules = [
    ("1", "Kunden bestämmer VAD", "AI-assistenten föreslår och utför, men kunden äger prioritering och acceptans."),
    ("2", "Allt dokumenteras", "Ingen kod utan wiki-sida. Alla beslut, kommandon och konfigurationer sparas."),
    ("3", "Godkänn först, avsluta sen", "Ingen task markeras Klart före kund-godkännande."),
    ("4", "Arbeta i cykler", "Bygg en liten bit, testa, godkänn. Inga stora hopp."),
    ("5", "Håll Kanban uppdaterad", "AI-assistenten ansvarar. Kunden ska alltid kunna se status."),
    ("6", "Två GitHub-repos", "Wiki (rfid-manager) och kod (rfid-manager-android) separeras."),
    ("7", "master = default branch", "Inte main. Alla repos följer detta."),
    ("8", "Fråga om osäker", "AI-assistenten frågar hellre än gissar. Kunden svarar."),
]

for i, (num, title, desc) in enumerate(rules):
    col = i % 2
    row = i // 2
    l = 0.5 + col * 6.3
    t = 1.6 + row * 1.4
    card(s, l, t, 5.9, 1.2, f"Regel {num}: {title}", f"{desc}", ACCENT)

# === SLIDE 10: End ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
tbox(s, 0.8, 2.0, 11.7, 1.0, "Sammanfattning", size=42, color=WHITE, bold=True)
tbox(s, 0.8, 3.2, 11.7, 0.6, "AI-assisterad, kundstyrd, dokumenterad — en skalbar arbetsmetod", size=22, color=ACCENT)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.0), Inches(3), Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
tbox(s, 0.8, 4.3, 11.7, 0.5, "Joakim • 2026-06-14", size=18, color=GRAY)

output_path = "/home/joakim/llm-wiki/RFID-Manager-Arbetssatt.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
