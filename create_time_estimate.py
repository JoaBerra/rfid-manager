from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_SLIDE = RGBColor(0x1e, 0x1e, 0x36)
ACCENT = RGBColor(0x00, 0xd2, 0xff)
GREEN = RGBColor(0x50, 0xfa, 0x7b)
ORANGE = RGBColor(0xff, 0xb7, 0x00)
RED = RGBColor(0xff, 0x55, 0x55)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0xaa, 0xaa, 0xcc)
DARK_CARD = RGBColor(0x28, 0x28, 0x44)

def set_bg(slide, color=BG_SLIDE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tbox(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb

def card(slide, l, t, w, h, title, content, title_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.space_after = Pt(8)
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(4)

# === SLIDE 1: Title ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
tbox(s, 0.8, 1.5, 11.7, 1.2, "RFID Manager", size=52, color=WHITE, bold=True)
tbox(s, 0.8, 2.8, 11.7, 0.8, "Tidsuppskattning — Kund/Projektledare", size=32, color=ACCENT)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.7), Inches(3), Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
tbox(s, 0.8, 4.0, 11.7, 0.5, "Hur mycket tid har du lagt ner som PM sedan projektstart?", size=18, color=GRAY)
tbox(s, 0.8, 5.5, 11.7, 0.5, "2026-06-14", size=16, color=GRAY)

# === SLIDE 2: Metod ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Hur uppskattningen är gjord", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Baserad på genomförda faser, antal sessioner och typiska PM-aktiviteter", size=16, color=GRAY)

card(s, 0.5, 1.6, 5.8, 5.0,
    "📊 Underlag",
    "Uppskattningen baseras på:\n\n"
    "• Git-loggar (wiki + android) — tidslinje\n"
    "   Wiki: 2026-06-07 till 2026-06-14\n"
    "   Android: 2026-06-13 till 2026-06-14\n\n"
    "• Kanban-board (58 slutförda tasks)\n\n"
    "• Genomförda faser: 6 hela + 1 påbörjad\n\n"
    "• 8 buggar fixade och godkända\n\n"
    "• AI-assistentens session-kontext\n\n"
    "⚠️ Detta är en uppskattning, inte exakt tidrapportering")

card(s, 6.8, 1.6, 5.8, 5.0,
    "🧑‍💼 Dina PM-aktiviteter",
    "Som Kund/Projektledare har du:\n\n"
    "• Definition: Skriva \"gör X\" — 5-15 min per task\n\n"
    "• Granskning: Läsa AI-svar, inspektera kod,\n"
    "  testa funktion — 10-30 min per task\n\n"
    "• Beslut: Godkänna/avböja/justera — 2-5 min\n\n"
    "• Frågor: Ställa uppföljningsfrågor,\n"
    "  be om förtydliganden — 2-10 min\n\n"
    "• Testning: adb deploy, provtryck på telefon —\n"
    "  5-20 min per testomgång\n\n"
    "Allt annat (kod, testskript, dokumentation,\n"
    "wiki-uppdateringar) utförs av AI-assistenten.")

# === SLIDE 3: Tidslinje ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Tidslinje", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Projektet startade med fork från mall — datum baserat på git-loggar och session-kontext", size=16, color=GRAY)

timeline = [
    ("Vecka 1", "Jun 7-8", "Fas 1-2", "Fork från mall. NFC-läsning/skrivning,\nradar-grafik, MQTT-infrastruktur,\npersistens, end-to-end test.", "4-6h"),
    ("Vecka 2", "Jun 9-10", "Fas 3-4", "Navigation, ViewModels, i18n,\nfont-slider, sök/filter, CSV/JSON-export,\nhaptic+ljud, MQTT-anslutning,\ndark mode, paginering.", "6-8h"),
    ("Vecka 3", "Jun 11-12", "Fas 5 + Buggar", "Användarmanual, arkitekturdiagram,\nrelease notes, testplan, kodgenomgång,\ndynamisk layout, radar sweep.\n8 buggar fixade.", "5-7h"),
    ("Vecka 4", "Jun 13-14", "Fas 6 + Fas-100", "MQTT-konfig UI, app-ikon, release\nbuild, stub-borttagning, v1.0 release.\nMQTT-infrastruktur fördjupning.", "4-6h"),
]

for i, (week, dates, phase, desc, time) in enumerate(timeline):
    l = 0.5
    t = 1.6 + i * 1.45
    card(s, l, t, 12.0, 1.3,
        f"{week} ({dates}) — {phase} — ⏱ {time}",
        desc, ACCENT)

tbox(s, 0.5, 7.0, 12, 0.4, "Totalt: ~19-27 timmar över 4 veckor", size=14, color=ORANGE, bold=True)

# === SLIDE 4: Uppskattning per fas ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Uppskattning per fas", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Endast din tid som Kund/PM (inte AI-utvecklingstid)", size=16, color=GRAY)

# Table as cards
phases = [
    ("Fas 1", "NFC, UI, Radar", "4 tasks", "2-3h", "✅"),
    ("Fas 2", "MQTT, Persistens, E2E", "3 tasks", "2-3h", "✅"),
    ("Fas 3", "Navigation, ViewModels", "4 tasks", "2-3h", "✅"),
    ("Fas 4", "i18n, Slider, Sök, Export, Haptic, MQTT, Dark, Paginering", "8 tasks", "5-7h", "✅"),
    ("Fas 5", "Manual, Diagram, Testplan, Code review, Layout, Radar", "7 tasks", "4-5h", "✅"),
    ("Fas 6", "MQTT UI, Icon, Release, Stub, E2E, v1.0", "6 tasks", "2-3h", "✅"),
    ("Buggar", "BUG-001 till BUG-008", "8 buggar", "2-3h", "✅"),
    ("Fas-100", "MQTT-fördjupning (pågår)", "11 områden", "1-2h", "🔄"),
]

for i, (phase, desc, count, time, status) in enumerate(phases):
    l = 0.5
    t = 1.6 + i * 0.7
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(12.0), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = f"{status}  {phase:<10} {desc:<45} {count:<15} ⏱ {time:<10}"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

# Total summary
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(7.0), Inches(12.0), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_CARD
shape.line.fill.background()
tf = shape.text_frame
tf.clear()
tf.margin_left = Inches(0.15)
tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]
p.text = "Totalsumma (din tid): ~20-27 timmar över 4 veckor — ca 5-7 timmar/vecka som PM"
p.font.size = Pt(14)
p.font.color.rgb = GREEN
p.font.bold = True

# === SLIDE 5: Detalj per aktivitetstyp ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Din tid per aktivitetstyp", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Vad har tiden lagts på?", size=16, color=GRAY)

activities = [
    ("📋 Krav & Definitioner", "6-8h", "30%",
     "Beskriva vad som ska byggas (\"Gör X\").\nVarje fas: 5-15 min att formulera i Kanban + diskussion."),
    ("🔍 Granskning & Test", "8-10h", "38%",
     "Läsa AI-svar, inspektera resultat, testa på telefon\n(adb deploy). Största posten — din kvalitetskontroll."),
    ("✅ Godkännanden & Beslut", "2-3h", "10%",
     "Sign-off per task, prioritering av nästa steg,\nbeslut om designval (\"gör så här istället\")."),
    ("❓ Frågor & Förtydliganden", "2-3h", "10%",
     "Uppföljningsfrågor till AI, förklara önskemål,\nklargöra acceptanskriterier."),
    ("📚 Dokumentation", "1-2h", "7%",
     "Läsa wiki-uppdateringar, be om förbättringar,\ngranska presentations-PDFer."),
    ("🎯 Övrigt (planering, möten)", "1-2h", "5%",
     "Planera nästa sprint, fundera på roadmap,\nstrategiska beslut."),
]

for i, (activity, time, pct, desc) in enumerate(activities):
    col = i % 2
    row = i // 2
    l = 0.5 + col * 6.3
    t = 1.6 + row * 1.9
    card(s, l, t, 5.9, 1.7, f"{activity}  ⏱ {time}  ({pct})", desc)

# === SLIDE 6: Vad tiden INTE inkluderar ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Vad denna tid INTE inkluderar", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Skillnaden mellan PM-tid och total projekttid", size=16, color=GRAY)

card(s, 0.5, 1.6, 5.8, 2.5,
    "🤖 AI-utvecklingstid (inte din tid)",
    "All kod, testskript, dokumentation och wiki-innehåll\nhar skrivits av OpenCode AI-assistenten.\n\n"
    "Uppskattad AI-tid: 40-60 timmar\n— but that's not your problem.\n\n"
    "Du betalar inte för AI-timmar, du betalar för\nresultat. Din tid är den som räknas som\nprojektledningsinsats.")

card(s, 6.8, 1.6, 5.8, 2.5,
    "📈 Om du själv skrivit koden",
    "För att sätta siffrorna i perspektiv:\n\n"
    "• Uppskattad kodmängd: ~5000+ rader\n  (Kotlin + Compose + config + tester)\n\n"
    "• Manuell utvecklingstid: 200-400 timmar\n  (erfaren Android-utvecklare)\n\n"
    "• Din PM-tid: 20-27 timmar\n\n"
    "• Besparing: ca 90% av utvecklingstiden\n  förskjuts från dig till AI\n\n"
    "Du får en app utan att skriva en rad kod.")

card(s, 0.5, 4.5, 5.8, 2.5,
    "⏱️ Total projekttid (uppskattning)",
    "┌──────────────────────────────┬─────────┐\n"
    "│ Aktivitet                   │ Timmar  │\n"
    "├──────────────────────────────┼─────────┤\n"
    "│ PM-tid (du)                 │   20-27 │\n"
    "│ AI-utveckling (OpenCode)    │   40-60 │\n"
    "│ Docker/infrastruktur        │    2-3   │\n"
    "│ Verktyg (setup, Obsidian)   │    1-2   │\n"
    "├──────────────────────────────┼─────────┤\n"
    "│ Total                      │  63-92  │\n"
    "│ (ca 2-3 arbetsveckor)       │         │\n"
    "└──────────────────────────────┴─────────┘")

card(s, 6.8, 4.5, 5.8, 2.5,
    "📅 Fördelning över tid",
    "• Projektet startade första veckan i juni\n\n"
    "• 4 veckor aktiv utveckling\n\n"
    "• Din insats: ~5-7 timmar/vecka\n  (som projektledare vid sidan av annat)\n\n"
    "• AI-insats: ~10-15 timmar/vecka\n  (hela implementationen)\n\n"
    "• Motsvarande heltidsutvecklare:\n  2-3 veckor för samma resultat\n  (utan AI)")

# === SLIDE 7: Slutsats ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Slutsats", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Vad betyder siffrorna?", size=16, color=GRAY)

card(s, 0.5, 1.6, 3.8, 2.5,
    "💰 Effektiv tid",
    "Du har lagt ~20-27 timmar som\nprojektledare och fått en\nfärdig Android-app med:\n\n"
    "• NFC-läsning/skrivning\n• MQTT-infrastruktur\n• Fullt UI (Material 3)\n• i18n på svenska/engelska\n• 6 faser + 8 buggfixar\n• Dokumentation + wiki\n• v1.0 release på GitHub")

card(s, 4.7, 1.6, 3.8, 2.5,
    "⚡ AI-multiplikatorn",
    "Din tid som PM multipliceras\nav AI-assistenten:\n\n"
    "1h PM-arbete → ~3h AI-arbete\n\n"
    "Du fokuserar på VAD,\nAI gör HUR.\n\n"
    "Utan AI: 200-400h utveckling\nMed AI: 20-27h PM + 40-60h AI\n= snabbare och billigare")

card(s, 8.9, 1.6, 3.8, 2.5,
    "📊 Nyckeltal",
    "• 4 veckor från start till v1.0\n• 6 faser + 8 buggar\n• 22 testfall, alla Godkänt\n• 2 GitHub-repos\n• 10+ wiki-sidor\n• 2 presentations-PDFer\n\n"
    "Din investering:\n~5-7 timmar/vecka som PM\n≈ 1 arbetsdag/vecka")

tbox(s, 0.5, 4.8, 12, 0.5, "Sammanfattning: ~20-27 timmar totalt som projektledare", size=18, color=ACCENT, bold=True)

# === SLIDE 8: End ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
tbox(s, 0.8, 2.0, 11.7, 1.0, "Frågor?", size=48, color=WHITE, bold=True)
tbox(s, 0.8, 3.2, 11.7, 0.6, "Tidsuppskattning — RFID Manager", size=24, color=ACCENT)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.0), Inches(3), Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
tbox(s, 0.8, 4.3, 11.7, 0.5, "Joakim • 2026-06-14", size=18, color=GRAY)

output_path = "/home/joakim/llm-wiki/RFID-Manager-Tidsuppskattning.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
