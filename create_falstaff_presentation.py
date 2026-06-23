from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_SLIDE = RGBColor(0x1e, 0x1e, 0x36)
ACCENT = RGBColor(0x00, 0xd2, 0xff)
ACCENT2 = RGBColor(0x7c, 0x3a, 0xed)
GREEN = RGBColor(0x50, 0xfa, 0x7b)
ORANGE = RGBColor(0xff, 0xb7, 0x00)
RED = RGBColor(0xff, 0x55, 0x55)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0xaa, 0xaa, 0xcc)
DARK_CARD = RGBColor(0x28, 0x28, 0x44)

def set_bg(slide, color=BG_SLIDE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tbox(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb

def add_bullet_text(tf, text, size=13, color=GRAY, bold=False, space_before=4):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def card(slide, l, t, w, h, title, content, title_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.space_after = Pt(8)
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(4)
    return shape

def card_multiline(slide, l, t, w, h, title, lines, title_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.space_after = Pt(6)
    for line in lines:
        add_bullet_text(tf, line)
    return shape

# === SLIDE 1: Title ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
tbox(s, 0.8, 1.2, 11.7, 1.2, "RFID Manager", size=52, color=WHITE, bold=True)
tbox(s, 0.8, 2.5, 11.7, 0.8, "Utvecklingsmiljö — Falstaff", size=32, color=ACCENT)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.5), Inches(3), Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
tbox(s, 0.8, 3.8, 11.7, 0.5, "Organisation • Regler • Nomenklatur • Standarder • Processer • Verktyg • Testmiljöer", size=18, color=GRAY)
tbox(s, 0.8, 5.5, 11.7, 0.5, "2026-06-23", size=16, color=GRAY)

# === SLIDE 2: How everything connects ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Hur allt hänger ihop", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Från AI-assisterad utveckling till fysisk NFC-scanning — hela kedjan", size=16, color=GRAY)

card(s, 0.5, 1.6, 4.0, 2.5,
    "🤖 AI-assistent (OpenCode)",
    "• MCP-server → dashboard\n• Skriver kod, tester, wiki\n• Håller Kanban aktuell\n• Utför implementation\n• ~/.config/opencode/")

card(s, 4.8, 1.6, 4.0, 2.5,
    "📱 Android App",
    "• Kotlin + Jetpack Compose\n• NFC-scanning (eskortminne)\n• Paho MQTT-klient\n• Wi-Fi: 192.168.50.112\n• GitHub: rfid-manager-android")

card(s, 9.1, 1.6, 3.8, 2.5,
    "🌐 MQTT Broker (Mosquitto)",
    "• Docker-container\n• Port 1883 (TCP)\n• IP: 192.168.50.107\n• allow_anonymous\n• eclipse-mosquitto:2")

card(s, 0.5, 4.5, 4.0, 2.5,
    "💾 Python Backend",
    "• Subscriber → SQLite\n• Dashboard (FastAPI/SSE)\n  → port 8001\n• .venv (paho-mqtt, fastapi)\n• data/rfid_readings.db")

card(s, 4.8, 4.5, 4.0, 2.5,
    "📓 Wiki (Obsidian + GitHub)",
    "• Karpathy-pattern\n• index.md → specialsidor\n• Kanban-board\n• Produkt-Roadmap\n• All historik i log.md")

card(s, 9.1, 4.5, 3.8, 2.5,
    "🖥️ Utvecklingsmiljö",
    "• Ubuntu 24.04 (Noble)\n• Android Studio (Snap)\n• ADB (Wi-Fi)\n• 3 repos under\n  ~/projects/rfid/\n• Allt versionshanterat")

# === SLIDE 3: IP Categories ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Våra Intellectual Properties", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Sju kategorier som definierar vår utvecklingsmiljö", size=16, color=GRAY)

ip_categories = [
    ("🏛️  Organisation", ACCENT,
     "Roller: Kund/Projektledare + AI-assistent\nTvå GitHub-repos: rfid-manager (wiki), rfid-manager-android (kod)\nTre lokala kataloger under ~/projects/rfid/"),
    ("📜  Regler", GREEN,
     "Kunden bestämmer VAD, AI utför HUR\nAllt dokumenteras, inget godkänns före Kund\nmaster = default branch\nFråga hellre än gissa"),
    ("📖  Nomenklatur", ORANGE,
     "Topic-struktur: rfidmanager/<uid>/telemetry\nPackage: com.joakim.rfidmanager\nWiki-filer: Stor bokstav, bindestreck\nRelease-taggar: fas-X-sign-off-YYYY-MM-DD"),
    ("🏷️  Standarder", ACCENT2,
     "MQTT 3.1.1, QoS 1, JSON-payload\nKotlin + Jetpack Compose\nMaterial 3 design\nKarpathy LLM Wiki-pattern"),
    ("🔄  Processer", ACCENT,
     "Behov → Kanban → AI impl. → Kund testar → Godkänt?\nJa → Klart + Dokumentation\nNej → Justering\nRelease: feature-fas → tagg → GitHub Release"),
    ("🛠️  Verktyg", GREEN,
     "OpenCode (CLI AI), Obsidian (wiki)\nNeovim (editor), GitHub (git)\nDocker (miljöer), ADB (debug)\npython-pptx (presentationer)"),
    ("🧪  Testmiljöer", ORANGE,
     "Mosquitto Docker-container (port 1883)\nMQTT Explorer (GUI, ✅ installerad)\nPython subscriber + SQLite-databas\nDashboard (FastAPI) + MCP-server"),
]

for i, (title, color, content) in enumerate(ip_categories):
    col = i % 4
    row = i // 4
    l = 0.5 + col * 3.2
    t = 1.6 + row * 2.6
    card(s, l, t, 2.9, 2.3, title, content, color)

# === SLIDE 4: Organisation ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "🏛️  Organisation", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Roller, repositorier och katalogstruktur", size=16, color=GRAY)

card(s, 0.5, 1.6, 3.8, 2.5,
    "👤 Kund / Projektledare",
    "• Äger produkten\n• Definierar behov & krav\n• Prioriterar & granskar\n• Godkänner alla leveranser\n• Signerar acceptans")

card(s, 4.7, 1.6, 3.8, 2.5,
    "🤖 AI-assistent (OpenCode)",
    "• Utför implementation\n• Skriver dokumentation\n• Föreslår lösningar\n• Håller Kanban uppdaterad\n• MCP-server för dashboard")

card(s, 8.9, 1.6, 3.8, 2.5,
    "📂 ~/projects/rfid/",
    "• llm-wiki/ — hjärnan (wiki + kod)\n• rfid-manager-android/ — Android-app\n\nAlla tre är git-repos.\nllm-wiki och rfid-manager-android\npushas till GitHub.")

card(s, 0.5, 4.5, 6.0, 2.5,
    "🐙 GitHub Repos & Branches",
    "rfid-manager (github.com/JoaBerra/rfid-manager)\n  • Wiki + scripts + presentations + Android-kopia\n  • Default branch: main\n  • Innehåller wiki/, test/, dashboard/, setup/, mcp-server/\n\nrfid-manager-android (github.com/JoaBerra/rfid-manager-android)\n  • Android-app: Kotlin + Jetpack Compose\n  • Default branch: master\n  • Innehåller app/src/, build.gradle.kts, gradle/")

card(s, 6.8, 4.5, 5.8, 2.5,
    "🗺️ Katalogstruktur (llm-wiki)",
    "wiki/         — Obsidian-wiki (index, Kanban, arkitektur)\ntest/         — Python testskript (MQTT subscriber)\ndashboard/    — FastAPI/SSE dashboard\nsetup/        — bootstrap.sh, udev-rules\nmcp-server/   — MCP-server för OpenCode\ndata/         — SQLite-databas (rfid_readings.db)\nRFIDManager/  — Gammal Android-kopia\npresentations/— Färdiga PPTX-filer", ACCENT)

# === SLIDE 5: Regler & Standarder ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Regler & Standarder", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Spelregler, policies och tekniska standarder", size=16, color=GRAY)

card(s, 0.5, 1.6, 6.0, 2.5,
    "📜 Processregler (AGENTS.md)",
    "1. Kunden bestämmer VAD, AI föreslår HUR\n2. Allt dokumenteras i wikin\n3. Godkänn först, avsluta sen\n4. Arbeta i korta cykler\n5. Håll Kanban uppdaterad\n6. Två GitHub-repos\n7. master = default branch (android)\n8. Fråga om osäker")

card(s, 6.8, 1.6, 5.8, 2.5,
    "🏷️ Tekniska standarder",
    "MQTT 3.1.1 — publish/subscribe\nQoS 1 — at least once delivery\nJSON — lightweight payload-format\nSparkplug B — industristandard (planerad)\n\nKotlin + Jetpack Compose — Android\nMaterial 3 — designlanguage\nRoom/SQLite — lokal persistens\n\nKarpathy LLM Wiki — dokumentationspattern\npython-pptx — presentationsgenerering\nMarkdown — all wiki-innehåll")

card(s, 0.5, 4.5, 6.0, 2.5,
    "🌐 Nätverksstandarder",
    "Subnet: 192.168.50.0/24\nWired (eno1): 192.168.50.107 — broker, dashboard, subscriber\nWi-Fi: 192.168.50.112 — Android-telefon\nPort 1883 — MQTT (cleartext, lokalt nät)\nPort 8001 — Dashboard (FastAPI/uvicorn)\nPort 5000 — MCP-server (OpenCode)\nADB: 192.168.50.112:45979 (Wi-Fi)",
    ORANGE)

card(s, 6.8, 4.5, 5.8, 2.5,
    "📖 Nomenklatur & Konventioner",
    "Topics: rfidmanager/<UID>/telemetry\nGit commits: <action> | <vad> (<repo>)\nWiki-filer: Stor-bokstav-med-bindestreck.md\nRelease-tags: fas-X-sign-off-YYYY-MM-DD\nPackage: com.joakim.rfidmanager\nPython-skript: snake_case\nKotlin: camelCase (variabler), PascalCase (klasser)",
    GREEN)

# === SLIDE 6: Processer ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "🔄  Processer", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Från behov till leverans — hela utvecklingsflödet", size=16, color=GRAY)

flow = [
    ("🎯", "Behov", "Kund identifierar\nny funktion eller\nförbättring"),
    ("📋", "Kanban", "Läggs i Att göra.\nPrioriteras.\nAccepanskriterier."),
    ("🤖", "AI impl.", "OpenCode utför:\nkod, tester,\ndokumentation"),
    ("👤", "Kund testar", "Live-granskning.\nTestkörning.\nInspektion."),
    ("✅", "Godkänt?", "Ja → Klart + sign-off\nNej → justering,\nny granskning"),
    ("📝", "Dokument.", "Wiki uppdateras.\nlog.md, Kanban.\nBeslut sparas."),
]

for i, (emoji, title, desc) in enumerate(flow):
    l = 0.5 + i * 2.1
    card(s, l, 1.6, 1.9, 2.8, f"{emoji} {title}", desc, ACCENT)

card(s, 0.5, 4.8, 12.0, 2.3,
    "📊 Planeringshierarki",
    "┌─────────────────────────────────────────────────────────────────────────────────────┐\n"
    "│                                                                                     │\n"
    "│  Produkt-Roadmap — långsiktig vision (faser 1-6+)                                   │\n"
    "│       ↓                                                                             │\n"
    "│  Kanban — veckovis taktik (Att göra → Pågår → Klart)                                │\n"
    "│       ↓                                                                             │\n"
    "│  Session log — daglig aktivitet (log.md, AGENTS.md)                                  │\n"
    "│                                                                                     │\n"
    "│  Release-process: Feature-fas klar → Kund godkänner → tagg → GitHub Release         │\n"
    "│  Bugghantering: BUG-XXX → Kanban → Fix → Kund testar → Godkänt → Stäng              │\n"
    "└─────────────────────────────────────────────────────────────────────────────────────┘",
    GREEN)

# === SLIDE 7: Verktyg ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "🛠️  Verktyg", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Hela stacken — från AI till fysisk enhet", size=16, color=GRAY)

tools = [
    ("🧠 OpenCode", "AI-assistent", "CLI-baserad.\nMCP-serverintegration.\nAgenter för parallella\ntasks."),
    ("📓 Obsidian", "Wiki/Kunskapsbas", "Karpathy-pattern.\nlänkar.\nKanban-plugin.\nbetter-mermaid."),
    ("💻 Neovim", "Editor", "Primär kod-editor.\nTerminalbaserad.\nAI via OpenCode."),
    ("🐙 GitHub", "Versionshantering", "2 repos.\nGitHub Releases.\nTags per fas."),
    ("🐳 Docker", "Miljö", "Mosquitto broker.\nIsolerade tjänster.\nPort 1883, 9001."),
    ("📱 Android Studio", "App-utveckling", "Snap-installerad.\nSDK 36, build-tools 36/37.\nADB (Wi-Fi)."),
    ("🐍 Python", "Backend/Scripts", "paho-mqtt, fastapi,\nuvicorn, python-pptx.\n.venv per komponent."),
    ("🔧 ADB", "Debug", "Wi-Fi: 192.168.50.112:45979.\nUdev-regler för USB.\nSamsung + Google."),
    ("📊 MQTT Explorer", "Övervakning", "GUI för topics.\nRealtidsträd.\nFelsökning."),
]

for i, (emoji, title, desc) in enumerate(tools):
    col = i % 3
    row = i // 3
    l = 0.5 + col * 4.2
    t = 1.6 + row * 1.9
    card(s, l, t, 3.8, 1.6, f"{emoji} {title}", desc)

# === SLIDE 8: Testmiljöer ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "🧪  Testmiljöer", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Fysiska enheter, containrar och databaser", size=16, color=GRAY)

card(s, 0.5, 1.6, 4.0, 2.5,
    "📱 Fysiska enheter",
    "Samsung Galaxy Note 10\n  • NFC (13.56 MHz HF)\n  • Android 12\n  • Wi-Fi: 192.168.50.112\n  • ADB ansluten\n\nTestläge: daglig E2E-scanning\nvia MQTT till dashboard")

card(s, 4.8, 1.6, 4.0, 2.5,
    "🐳 Docker-containrar",
    "Mosquitto MQTT Broker\n  • eclipse-mosquitto:2\n  • Port 1883 (MQTT)\n  • Port 9001 (WS)\n  • allow_anonymous: true\n  • Data sparas i containern\n\nStart: docker start rfid-mqtt")

card(s, 9.1, 1.6, 3.8, 2.5,
    "🐍 Python-testmiljö",
    "Subscriber: test_subscriber_persist.py\n  • SQLite-persistens\n  • Färgkodad output\n  • --uid-filter\n\nDashboard: FastAPI/SSE\n  • Port 8001\n  • Realtidsuppdatering\n  • MCP-server integration\n\nSimulator: simulate_mobile_publish.py\n  • Testa utan telefon",
    GREEN)

card(s, 0.5, 4.5, 4.0, 2.5,
    "🗄️ Databaser",
    "SQLite: ~/projects/rfid/llm-wiki/data/rfid_readings.db\n  • Alla MQTT-meddelanden\n  • Enkel SQL-query\n  • Ingen server krävs\n\nApp: Room DB (lokal cache)\n  • SQLite på Android\n  • Synk via MQTT")

card(s, 4.8, 4.5, 4.0, 2.5,
    "🔍 Verifieringsverktyg",
    "mosquitto_sub — CLI-övervakning\ndocker exec rfid-mqtt mosquitto_sub -t \"#\" -v\n\nMQTT Explorer — GUI (✅ installerad)\ndocker logs rfid-mqtt — broker-loggar\n\nWireshark — paketanalys (✅ installerad)\nADB logcat — Android-loggar")

card(s, 9.1, 4.5, 3.8, 2.5,
    "🧪 Testplan (22 testfall)",
    "Kategorier:\n  NAV — Navigation (3)\n  SCAN — Scan & NFC (7)\n  READ — Readings & Data (6)\n  SETT — Settings & Config (4)\n  E2E — End-to-End (2)\n\nAlla 22 Godkända ✅",
    ORANGE)

# === SLIDE 9: MQTT Data Flow ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "MQTT — Dataflöde", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Från NFC-scanning till persisterad data i SQLite", size=16, color=GRAY)

flow_steps = [
    ("1️⃣", "NFC Scan", "Telefon läser\nRFID-tagg\n→ extraherar UID\n→ bygger JSON"),
    ("2️⃣", "Publish MQTT", "App publicerar till\nrfidmanager/<uid>/telemetry\nQoS 1"),
    ("3️⃣", "Broker Router", "Mosquitto tar emot\noch vidarebefordrar\ntill alla subscribers"),
    ("4️⃣", "Python Subscriber", "Tar emot → parsar\nJSON → loggar till\nSQLite-databas"),
    ("5️⃣", "Dashboard", "FastAPI/SSE visar\nrealtidsdata.\n/api/stats, /api/messages"),
    ("6️⃣", "MCP-server", "OpenCode kan fråga:\n\"Vad är senaste\nstatus?\" via 4 tools"),
]

for i, (emoji, title, desc) in enumerate(flow_steps):
    l = 0.5 + i * 2.1
    card(s, l, 1.6, 1.9, 2.8, f"{emoji} {title}", desc, ACCENT)

card(s, 0.5, 4.8, 12.0, 2.3,
    "🔗 Payload-exempel (JSON)",
    '{\n'
    '  "uid": "047B05CA885884",\n'
    '  "type": "rfid",\n'
    '  "timestamp": 1718300000000,\n'
    '  "dataPreview": "042B...",\n'
    '  "memoryBank": "EPC",\n'
    '  "address": 2,\n'
    '  "length": 6,\n'
    '  "source": "scan",\n'
    '  "sparkplugJson": false\n'
    '}\n\n'
    'Topic: rfidmanager/047B05CA885884/telemetry\n'
    'Broker: 192.168.50.107:1883 | QoS 1 | cleanSession: true',
    GREEN)

# === SLIDE 10: Dashboard & MCP ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Dashboard & MCP-server", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Realtidsvisualisering och AI-assistentintegration", size=16, color=GRAY)

card(s, 0.5, 1.6, 6.0, 2.5,
    "📊 Web Dashboard (FastAPI + SSE)",
    "Stack: FastAPI, uvicorn, paho-mqtt\n\nEndpoints:\n  GET  /api/stats      — total messages, unique UIDs, broker status\n  GET  /api/messages   — senaste meddelandena\n  GET  /api/live       — SSE-ström för realtid\n\nFrontend: statisk HTML med mörkt tema\n  • Live-flöde med auto-scroll\n  • Statistik-kort\n  • Detaljvy per meddelande\n\nStart: cd dashboard && .venv/bin/uvicorn app.main:app --host 0.0.0.0 --port 8001",
    ACCENT)

card(s, 6.8, 1.6, 5.8, 2.5,
    "🤖 MCP-server (OpenCode)",
    "Stack: Python MCP SDK, httpx, paho-mqtt\n\nTools:\n  get_stats       — dashboard-statistik\n  get_messages    — senaste meddelanden\n  publish_mqtt    — publicera till broker\n  get_live_events — strömma live-data i 5s\n\nKonfig: ~/.config/opencode/opencode.jsonc\n  Ansluter till dashboard (8001) och broker (1883)\n\nAI-assistenten kan fråga:\n  \"Hur många meddelanden har vi fått?\"\n  \"Skicka ett testmeddelande\"\n  \"Vad är senaste status?\"",
    GREEN)

card(s, 0.5, 4.5, 6.0, 2.5,
    "🏗️ Arkitektur",
    "┌────────────────┐     ┌──────────────────┐     ┌──────────────────┐\n"
    "│  Android App   │────→│  MQTT Broker     │────→│  Python Subscriber│\n"
    "│  (Paho Client) │     │  (Mosquitto)     │     │  → SQLite.db     │\n"
    "└────────────────┘     └──────┬───────────┘     └──────────────────┘\n"
    "                              │                                    \n"
    "                              ▼                                    \n"
    "                    ┌────────────────────┐    ┌──────────────────┐\n"
    "                    │  Web Dashboard     │←───│  MCP-server      │\n"
    "                    │  (FastAPI + SSE)   │    │  (OpenCode tools)│\n"
    "                    └────────────────────┘    └──────────────────┘\n"
    "                                                      │\n"
    "                                                      ▼\n"
    "                                            ┌──────────────────┐\n"
    "                                            │  OpenCode AI     │\n"
    "                                            │  (fråga/svara)   │\n"
    "                                            └──────────────────┘",
    ACCENT2)

card(s, 6.8, 4.5, 5.8, 2.5,
    "🚀 Starta allt — one-liners",
    "Broker:  docker start rfid-mqtt\n\nDashboard:  cd ~/projects/rfid/llm-wiki && \\\n  dashboard/.venv/bin/uvicorn \\\n  app.main:app --host 0.0.0.0 --port 8001\n\nSubscriber:  test/fas2-mqtt/mqtt/ \\\n  ./test_subscriber_persist.py\n\nMCP:  mcp-server/.venv/bin/python \\\n  mcp-server/server.py\n\nADB:  adb connect 192.168.50.112:45979\n\nTest:  docker exec rfid-mqtt mosquitto_sub \\\n  -h localhost -t \"rfidmanager/#\" -v",
    ORANGE)

# === SLIDE 11: Kom-igång-guide ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Kom i gång — checklista", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Steg för att sätta upp miljön från scratch", size=16, color=GRAY)

card(s, 0.5, 1.6, 6.0, 2.5,
    "🖥️ System (bootstrap.sh)",
    "1. Ubuntu 24.04 — bas OS\n2. Docker + docker-compose-v2\n   → docker start rfid-mqtt\n3. Python 3.12.3 (pre-installerat)\n4. Java 21 OpenJDK (pre-installerat)\n5. Android Studio (snap)\n   → sudo snap install android-studio --classic\n6. Android SDK (build-tools 36/37, android-36)\n7. ADB + udev-regler (51-android.rules)\n8. Git — klona repos",
    GREEN)

card(s, 6.8, 1.6, 5.8, 2.5,
    "📦 Projekt (efter bootstrap)",
    "1. mkdir -p ~/projects/rfid/\n2. Git-klona:\n   github.com/JoaBerra/rfid-manager\n     → ~/projects/rfid/llm-wiki\n   github.com/JoaBerra/rfid-manager-android\n     → ~/projects/rfid/rfid-manager-android\n3. Skapa Python .venv:\n   test/fas2-mqtt/.venv\n   dashboard/.venv\n   mcp-server/.venv\n4. Starta broker + dashboard\n5. Anslut ADB (Wi-Fi)\n6. Bygg app: ./gradlew assembleDebug\n7. Verifiera: skanna NFC-tagg",
    ORANGE)

card(s, 0.5, 4.5, 6.0, 2.5,
    "🔧 Konfiguration ny maskin",
    "~/.bashrc tillägg:\n  export ANDROID_HOME=$HOME/Android/Sdk\n  export PATH=$PATH:$ANDROID_HOME/platform-tools\n\ndocker-grupp (kräver logout/login):\n  sudo usermod -aG docker $USER\n\nOpenCode-konfig:\n  ~/.config/opencode/opencode.jsonc\n  med MCP-server entry\n\nObsidian: öppna wiki/ som vault\n  Installera plugins: kanban, mermaid",
    ACCENT)

card(s, 6.8, 4.5, 5.8, 2.5,
    "✅ Verifieringsrutiner (Fas-400)",
    "1. Docker fungerar? → docker ps\n2. Broker lyssnar? → mosquitto_sub -t \"#\"\n3. ADB ansluten? → adb devices\n4. Dashboard svarar? → curl localhost:8001/api/stats\n5. MCP-server aktiv? → /api/stats via OpenCode\n6. Subscriber skriver? → cat data/rfid_readings.db\n7. App bygger? → ./gradlew assembleDebug\n8. E2E: skanna → syns i dashboard?\n9. Wiki tillgänglig? → Obsidian öppnar wiki/",
    GREEN)

# === SLIDE 12: Milstolpar ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
tbox(s, 0.5, 0.3, 12, 0.8, "Milstolpar", size=36, color=WHITE, bold=True)
tbox(s, 0.5, 1.0, 12, 0.4, "Vad har vi byggt hittills?", size=16, color=GRAY)

milestones = [
    ("Fas 1", "2026-05-26", "Initial struktur, Android Studio, NFC-läsning", ACCENT),
    ("Fas 2", "2026-06-04", "Eskortminne-läs/skriv, MQTT, persistens, E2E", ACCENT),
    ("Fas 3", "2026-06-10", "Navigation, ViewModels, spacing, polish, PC-stöd", ACCENT),
    ("Fas 4", "2026-06-11", "i18n, font-slider, sök, export, haptik, MQTT, dark mode, paginering", ACCENT),
    ("Fas 5", "2026-06-13", "Manual, arkitekturdiagram, testplan, kodgenomgång, radar trail", ACCENT),
    ("Fas 6", "2026-06-13", "v1.0 Release, app-ikon, release-build, GitHub Release", ACCENT),
    ("Fas 200", "2026-06-19", "MQTT Realtidsdashboard (FastAPI + SSE)", GREEN),
    ("Fas 300", "2026-06-19", "MCP-server för OpenCode (4 tools)", GREEN),
    ("Fas 400", "2026-06-21", "Teknikmiljö validering + bootstrap.sh", GREEN),
    ("Fas 500", "2026-06-21", "Miljöflytt sixten → falstaff", GREEN),
]

for i, (fas, date, desc, color) in enumerate(milestones):
    col = i % 2
    row = i // 2
    l = 0.5 + col * 6.3
    t = 1.6 + row * 1.1
    card(s, l, t, 5.9, 0.9, f"{fas} ({date})", desc, color)

tbox(s, 0.5, 7.0, 12.0, 0.4, "Nästa: Fas-100 (MQTT-fördjupning) • Fas-101 (MQTT-klientkonfiguration i appen)", size=14, color=ORANGE, align=PP_ALIGN.CENTER)

# === SLIDE 13: Sammanfattning ===
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
tbox(s, 0.8, 1.5, 11.7, 1.0, "Sammanfattning", size=42, color=WHITE, bold=True)
tbox(s, 0.8, 2.8, 11.7, 0.6, "En komplett utvecklingsmiljö för RFID-applikationer på Android", size=22, color=ACCENT)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.6), Inches(3), Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

tbox(s, 0.8, 3.9, 11.7, 2.5,
    "🏛️ Organisation: Kund + AI-assistent, 2 GitHub-repos, 3 lokala kataloger\n"
    "📜 Regler: Kunden bestämmer, allt dokumenteras, godkänn före avslut\n"
    "📖 Nomenklatur: Konsekventa namnkonventioner i hela stacken\n"
    "🏷️ Standarder: MQTT 3.1.1, Kotlin/Compose, Material 3, Karpathy-wiki\n"
    "🔄 Processer: Korta cykler, Kanban, sign-off, GitHub Releases\n"
    "🛠️ Verktyg: OpenCode, Obsidian, Neovim, GitHub, Docker, Android Studio\n"
    "🧪 Testmiljöer: Fysisk enhet, Mosquitto, SQLite, Dashboard, MCP-server",
    size=16, color=GRAY)

tbox(s, 0.8, 5.8, 11.7, 0.5, "Joakim • 2026-06-23", size=18, color=GRAY)

# Save PPTX
output_pptx = "/home/joakim/projects/rfid/llm-wiki/RFID-Manager-Falstaff-Miljo.pptx"
prs.save(output_pptx)
print(f"Saved: {output_pptx}")
